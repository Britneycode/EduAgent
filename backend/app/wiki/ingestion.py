from __future__ import annotations

import asyncio
import hashlib
from io import BytesIO
import json
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.wiki import WikiEntry
from app.wiki.vector_store import ChromaHttpVectorStore, VectorStore

logger = logging.getLogger(__name__)

_MAX_CHUNK_CHARS = 1800
_MIN_CHUNK_CHARS = 180
_SUPPORTED_UPLOAD_SUFFIXES = {".md", ".markdown", ".txt", ".pdf", ".pptx"}


@dataclass(slots=True)
class DocumentChunk:
    """从 Markdown 解析出的知识块。"""

    chunk_id: str
    title: str
    content: str
    course_id: str = ""
    chapter: str = ""
    section: str = ""
    tags: list[str] = field(default_factory=list)
    content_type: str = "original"
    source_agent: str | None = None
    source_name: str = ""
    mime_type: str = ""


@dataclass(slots=True)
class UploadedDocumentIngestionResult:
    """上传资料入库结果。"""

    filename: str
    title: str
    course_id: str
    content_type: str
    chunk_count: int
    chunk_ids: list[str]
    char_count: int
    chapter: str
    section: str


class DocumentIngestionError(ValueError):
    """资料导入失败。"""


class UnsupportedDocumentTypeError(DocumentIngestionError):
    """不支持的资料格式。"""


class KnowledgeIngestion:
    """知识导入管道：Markdown → 分块 → 向量化 → 入库。"""

    def __init__(
        self,
        vector_store: VectorStore | ChromaHttpVectorStore,
        session: AsyncSession | None = None,
    ) -> None:
        self._vector_store = vector_store
        self._session = session

    async def ingest_course(
        self,
        knowledge_dir: str | Path,
        *,
        course_id: str = "",
    ) -> int:
        """导入整个课程知识库，返回导入的 chunk 总数。"""
        knowledge_path = Path(knowledge_dir)
        metadata_file = knowledge_path / "metadata.json"
        if not metadata_file.exists():
            logger.error("元数据文件不存在: %s", metadata_file)
            return 0

        with open(metadata_file, encoding="utf-8") as f:
            metadata = json.load(f)

        total_chunks = 0
        for index, chapter in enumerate(metadata.get("chapters", []), start=1):
            if not isinstance(chapter, dict):
                continue

            md_files = _resolve_chapter_files(knowledge_path, chapter, index)
            for md_file in md_files:
                if not md_file.exists():
                    logger.warning("章节文件不存在，跳过: %s", md_file)
                    continue

                chapter_for_file = dict(chapter)
                if len(md_files) > 1:
                    chapter_for_file.setdefault(
                        "source_id",
                        _build_source_id(knowledge_path, md_file),
                    )
                chunks = self._parse_markdown(
                    md_file,
                    chapter_for_file,
                    course_id=course_id,
                )
                if chunks:
                    await self._store_chunks(chunks)
                    total_chunks += len(chunks)
                    logger.info(
                        "导入章节 %s / %s: %d 个知识块",
                        chapter.get("title", ""),
                        md_file.name,
                        len(chunks),
                    )

        logger.info("知识导入完成，共 %d 个知识块", total_chunks)
        return total_chunks

    async def ingest_uploaded_document(
        self,
        *,
        filename: str,
        content: bytes,
        mime_type: str = "",
        chapter: str | None = None,
        section: str | None = None,
        tags: list[str] | None = None,
        course_id: str | None = None,
    ) -> UploadedDocumentIngestionResult:
        """导入用户上传的 Markdown/TXT/PDF/PPTX 资料。"""
        source_name = Path(filename).name or "uploaded-document"
        title = Path(source_name).stem or source_name
        text = extract_upload_text(filename=source_name, content=content)
        if not text.strip():
            raise DocumentIngestionError("上传资料未解析出可入库文本")

        resolved_chapter = (chapter or "uploaded").strip() or "uploaded"
        resolved_section = (section or "").strip()
        resolved_tags = [tag.strip() for tag in tags or [] if tag.strip()]
        chunks = build_uploaded_document_chunks(
            title=title,
            text=text,
            source_name=source_name,
            mime_type=mime_type,
            course_id=(course_id or "").strip(),
            chapter=resolved_chapter,
            section=resolved_section,
            tags=resolved_tags,
        )
        await self._store_chunks(chunks)
        return UploadedDocumentIngestionResult(
            filename=source_name,
            title=title,
            course_id=(course_id or "").strip(),
            content_type=_detect_upload_kind(source_name),
            chunk_count=len(chunks),
            chunk_ids=[chunk.chunk_id for chunk in chunks],
            char_count=len(text),
            chapter=resolved_chapter,
            section=resolved_section,
        )

    def _parse_markdown(
        self,
        md_file: Path,
        chapter_info: dict[str, Any],
        *,
        course_id: str = "",
    ) -> list[DocumentChunk]:
        """按三级标题 (###) 分割 Markdown 文件为知识块。"""
        text = md_file.read_text(encoding="utf-8")
        chapter_id = str(
            chapter_info.get("id") or chapter_info.get("chapter_id") or ""
        )
        source_id = str(chapter_info.get("source_id") or "").strip()
        if not chapter_id:
            logger.warning("章节缺少 id/chapter_id，跳过: %s", md_file)
            return []
        sections = chapter_info.get("sections", [])

        # 构建 section title -> section id 映射
        section_map: dict[str, str] = {}
        for sec in sections:
            section_map[sec["title"]] = sec["id"]

        chunks: list[DocumentChunk] = []
        current_section_id = ""

        # 按标题行分割
        lines = text.split("\n")
        current_title = ""
        current_content_lines: list[str] = []
        current_level = 0

        for line in lines:
            heading_match = re.match(r"^(#{1,3})\s+(.+)$", line)
            if heading_match:
                level = len(heading_match.group(1))
                title = heading_match.group(2).strip()

                # 先保存之前积累的块
                if current_level == 3 and current_title:
                    content = "\n".join(current_content_lines).strip()
                    if content:
                        chunk_id = _build_chunk_id(
                            course_id,
                            chapter_id,
                            current_title,
                            source_id=source_id,
                            ordinal=len(chunks) + 1,
                        )
                        chunks.append(
                            DocumentChunk(
                                chunk_id=chunk_id,
                                title=current_title,
                                content=content,
                                course_id=course_id,
                                chapter=chapter_id,
                                section=current_section_id,
                                source_name=md_file.name,
                            )
                        )

                if level == 2:
                    # 尝试匹配 section id（去掉编号前缀）
                    clean_title = re.sub(r"^\d+\.\d+\s*", "", title)
                    current_section_id = section_map.get(
                        clean_title, section_map.get(title, "")
                    )
                    current_title = ""
                    current_content_lines = []
                    current_level = 2
                elif level == 3:
                    current_title = title
                    current_content_lines = []
                    current_level = 3
                elif level == 1:
                    # 章标题，不作为 chunk
                    current_title = ""
                    current_content_lines = []
                    current_level = 1
            else:
                if current_level >= 2:
                    current_content_lines.append(line)

        # 保存最后一个块
        if current_level == 3 and current_title:
            content = "\n".join(current_content_lines).strip()
            if content:
                chunk_id = _build_chunk_id(
                    course_id,
                    chapter_id,
                    current_title,
                    source_id=source_id,
                    ordinal=len(chunks) + 1,
                )
                chunks.append(
                    DocumentChunk(
                        chunk_id=chunk_id,
                        title=current_title,
                        content=content,
                        course_id=course_id,
                        chapter=chapter_id,
                        section=current_section_id,
                        source_name=md_file.name,
                    )
                )

        # 如果没有三级标题，按二级标题分块
        if not chunks:
            chunks = self._parse_by_h2(
                text,
                chapter_id,
                section_map,
                course_id=course_id,
                source_id=source_id,
                source_name=md_file.name,
            )

        return chunks

    def _parse_by_h2(
        self,
        text: str,
        chapter_id: str,
        section_map: dict[str, str],
        *,
        course_id: str = "",
        source_id: str = "",
        source_name: str = "",
    ) -> list[DocumentChunk]:
        """按二级标题分块的后备方案。"""
        chunks: list[DocumentChunk] = []
        lines = text.split("\n")
        current_title = ""
        current_section_id = ""
        current_content_lines: list[str] = []

        for line in lines:
            heading_match = re.match(r"^##\s+(.+)$", line)
            if heading_match:
                if current_title and current_content_lines:
                    content = "\n".join(current_content_lines).strip()
                    if content:
                        chunk_id = _build_chunk_id(
                            course_id,
                            chapter_id,
                            current_title,
                            source_id=source_id,
                            ordinal=len(chunks) + 1,
                        )
                        chunks.append(
                            DocumentChunk(
                                chunk_id=chunk_id,
                                title=current_title,
                                content=content,
                                course_id=course_id,
                                chapter=chapter_id,
                                section=current_section_id,
                                source_name=source_name,
                            )
                        )
                title = heading_match.group(1).strip()
                clean_title = re.sub(r"^\d+\.\d+\s*", "", title)
                current_title = clean_title
                current_section_id = section_map.get(
                    clean_title, section_map.get(title, "")
                )
                current_content_lines = []
            elif not re.match(r"^#\s+", line):
                current_content_lines.append(line)

        if current_title and current_content_lines:
            content = "\n".join(current_content_lines).strip()
            if content:
                chunk_id = _build_chunk_id(
                    course_id,
                    chapter_id,
                    current_title,
                    source_id=source_id,
                    ordinal=len(chunks) + 1,
                )
                chunks.append(
                    DocumentChunk(
                        chunk_id=chunk_id,
                        title=current_title,
                        content=content,
                        course_id=course_id,
                        chapter=chapter_id,
                        section=current_section_id,
                        source_name=source_name,
                    )
                )

        return chunks

    async def _store_chunks(self, chunks: list[DocumentChunk]) -> None:
        """批量写入向量存储和数据库。"""
        chunk_ids = [c.chunk_id for c in chunks]
        documents = [f"{c.title}\n\n{c.content}" for c in chunks]
        metadatas = [
            {
                "course_id": c.course_id,
                "chapter": c.chapter,
                "section": c.section,
                "title": c.title,
                "content_type": c.content_type,
                "source_agent": c.source_agent or "",
                "source_name": c.source_name,
                "mime_type": c.mime_type,
                "tags": c.tags,
            }
            for c in chunks
        ]

        await asyncio.to_thread(
            self._vector_store.add,
            chunk_ids=chunk_ids,
            documents=documents,
            metadatas=metadatas,
        )

        if self._session is not None:
            for chunk in chunks:
                existing = await self._session.scalar(
                    select(WikiEntry).where(WikiEntry.chunk_id == chunk.chunk_id)
                )
                if existing is None:
                    entry = WikiEntry(
                        course_id=chunk.course_id or None,
                        chapter=chunk.chapter,
                        section=chunk.section,
                        title=chunk.title,
                        content=chunk.content,
                        content_type=chunk.content_type,
                        source_agent=chunk.source_agent,
                        chunk_id=chunk.chunk_id,
                        tags=chunk.tags,
                    )
                    self._session.add(entry)
                else:
                    existing.chapter = chunk.chapter
                    existing.section = chunk.section
                    existing.course_id = chunk.course_id or None
                    existing.title = chunk.title
                    existing.content = chunk.content
                    existing.content_type = chunk.content_type
                    existing.source_agent = chunk.source_agent
                    existing.tags = chunk.tags
            await self._session.commit()


def extract_upload_text(*, filename: str, content: bytes) -> str:
    """从上传资料中抽取纯文本。"""
    suffix = Path(filename).suffix.lower()
    if suffix not in _SUPPORTED_UPLOAD_SUFFIXES:
        supported = "、".join(sorted(_SUPPORTED_UPLOAD_SUFFIXES))
        raise UnsupportedDocumentTypeError(f"暂不支持 {suffix or '未知'} 格式，仅支持 {supported}")

    if suffix in {".md", ".markdown", ".txt"}:
        return _decode_text(content)
    if suffix == ".pdf":
        return _extract_pdf_text(content)
    if suffix == ".pptx":
        return _extract_pptx_text(content)

    raise UnsupportedDocumentTypeError("暂不支持该资料格式")


def _resolve_chapter_file(
    knowledge_path: Path,
    chapter: dict[str, Any],
    index: int,
) -> Path:
    file_name = chapter.get("file") or chapter.get("filename")
    if isinstance(file_name, str) and file_name.strip():
        return knowledge_path / file_name

    chapter_id = str(chapter.get("id") or chapter.get("chapter_id") or "")
    candidates = list(knowledge_path.glob("*.md")) + list(
        knowledge_path.glob("*.markdown")
    )
    for path in candidates:
        stem = path.stem.lower()
        if chapter_id and chapter_id.lower() in stem:
            return path
        if stem.startswith(f"chapter_{index:02d}_") or stem.startswith(
            f"chapter{index}_"
        ):
            return path

    return knowledge_path / f"chapter_{index:02d}.md"


def _resolve_chapter_files(
    knowledge_path: Path,
    chapter: dict[str, Any],
    index: int,
) -> list[Path]:
    raw_files = chapter.get("files")
    if isinstance(raw_files, list) and raw_files:
        files: list[Path] = []
        for raw_file in raw_files:
            if isinstance(raw_file, str):
                file_name = raw_file
            elif isinstance(raw_file, dict):
                file_name = str(raw_file.get("file") or raw_file.get("filename") or "")
            else:
                continue
            if file_name.strip():
                files.append(knowledge_path / file_name)
        return files

    return [_resolve_chapter_file(knowledge_path, chapter, index)]


def _build_source_id(knowledge_path: Path, md_file: Path) -> str:
    try:
        relative = md_file.relative_to(knowledge_path).with_suffix("").as_posix()
    except ValueError:
        relative = md_file.stem
    return hashlib.sha1(relative.encode("utf-8")).hexdigest()[:10]


def _build_chunk_id(
    course_id: str,
    chapter_id: str,
    title: str,
    *,
    source_id: str = "",
    ordinal: int | None = None,
) -> str:
    if source_id and ordinal is not None:
        base = f"{chapter_id}_{source_id}_{ordinal:03d}_{title}"
    elif source_id:
        base = f"{chapter_id}_{source_id}_{title}"
    else:
        base = f"{chapter_id}_{title}"
    return f"{course_id}:{base}" if course_id else base


def build_uploaded_document_chunks(
    *,
    title: str,
    text: str,
    source_name: str,
    mime_type: str,
    course_id: str,
    chapter: str,
    section: str,
    tags: list[str],
) -> list[DocumentChunk]:
    """把上传资料切分为可检索知识块。"""
    normalized_text = _normalize_text(text)
    sections = _split_uploaded_sections(normalized_text)
    digest = hashlib.sha1(
        f"{course_id}\n{source_name}\n{normalized_text}".encode("utf-8")
    ).hexdigest()[:12]
    chunks: list[DocumentChunk] = []

    for index, (section_title, section_text) in enumerate(sections, start=1):
        for part_index, part_text in enumerate(_split_long_text(section_text), start=1):
            chunk_title = section_title or title
            if len(sections) > 1 and chunk_title == title:
                chunk_title = f"{title} {index}"
            if part_index > 1:
                chunk_title = f"{chunk_title}（{part_index}）"
            chunks.append(
                DocumentChunk(
                    chunk_id=f"upload_{digest}_{len(chunks) + 1}",
                    title=chunk_title[:180],
                    content=part_text,
                    course_id=course_id,
                    chapter=chapter,
                    section=section,
                    tags=["uploaded", *tags],
                    content_type="uploaded_document",
                    source_agent="upload",
                    source_name=source_name,
                    mime_type=mime_type,
                )
            )

    if not chunks:
        raise DocumentIngestionError("上传资料内容过短，无法生成知识块")
    return chunks


def _detect_upload_kind(filename: str) -> str:
    suffix = Path(filename).suffix.lower().lstrip(".")
    if suffix in {"md", "markdown"}:
        return "markdown"
    return suffix or "text"


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="ignore")


def _extract_pdf_text(content: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DocumentIngestionError("缺少 pypdf 依赖，无法解析 PDF") from exc

    try:
        reader = PdfReader(BytesIO(content))
        pages = [page.extract_text() or "" for page in reader.pages]
    except Exception as exc:
        raise DocumentIngestionError("PDF 解析失败，请检查文件是否损坏") from exc
    return "\n\n".join(page.strip() for page in pages if page.strip())


def _extract_pptx_text(content: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentIngestionError("缺少 python-pptx 依赖，无法解析 PPTX") from exc

    try:
        presentation = Presentation(BytesIO(content))
    except Exception as exc:
        raise DocumentIngestionError("PPTX 解析失败，请检查文件是否损坏") from exc

    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
        if lines:
            slides.append(f"## Slide {index}\n" + "\n".join(lines))
    return "\n\n".join(slides)


def _normalize_text(text: str) -> str:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    return normalized.strip()


def _split_uploaded_sections(text: str) -> list[tuple[str, str]]:
    heading_pattern = re.compile(r"^(#{1,3})\s+(.+)$", re.MULTILINE)
    matches = list(heading_pattern.finditer(text))
    if not matches:
        return [("", text)]

    sections: list[tuple[str, str]] = []
    for idx, match in enumerate(matches):
        start = match.end()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else len(text)
        title = match.group(2).strip()
        content = text[start:end].strip()
        if content:
            sections.append((title, content))
    return sections or [("", text)]


def _split_long_text(text: str) -> list[str]:
    paragraphs = [item.strip() for item in re.split(r"\n\s*\n", text) if item.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        candidate = f"{current}\n\n{paragraph}".strip() if current else paragraph
        if len(candidate) <= _MAX_CHUNK_CHARS:
            current = candidate
            continue
        if current:
            chunks.append(current)
            current = ""
        if len(paragraph) <= _MAX_CHUNK_CHARS:
            current = paragraph
        else:
            chunks.extend(_split_by_sentence(paragraph))
    if current:
        chunks.append(current)
    return [chunk for chunk in chunks if len(chunk) >= _MIN_CHUNK_CHARS or len(chunks) == 1]


def _split_by_sentence(text: str) -> list[str]:
    pieces = re.split(r"(?<=[。！？!?；;])", text)
    chunks: list[str] = []
    current = ""
    for piece in pieces:
        if not piece:
            continue
        candidate = f"{current}{piece}"
        if len(candidate) <= _MAX_CHUNK_CHARS:
            current = candidate
            continue
        if current:
            chunks.append(current.strip())
        current = piece.strip()
    if current:
        chunks.append(current.strip())
    return chunks
